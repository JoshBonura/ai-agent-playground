from __future__ import annotations

from ...core.logging import get_logger

log = get_logger(__name__)
import io
import re
from fastapi import HTTPException
from .ocr import MissingTesseractError, _install_hint
from ...core.settings import SETTINGS
from .ocr import ocr_image_bytes

_WS_RE = re.compile("[ \\t]+")


def _squeeze(s: str) -> str:
    s = (s or "").replace("\xa0", " ")
    s = _WS_RE.sub(" ", s)
    return s.strip()


def _clip(s: str, limit: int) -> str:
    if limit > 0 and len(s) > limit:
        return s[:limit] + "…"
    return s


def _shape_text(shape) -> list[str]:
    out: list[str] = []
    if getattr(shape, "has_text_frame", False):
        for p in shape.text_frame.paragraphs:
            txt = _squeeze("".join(r.text for r in p.runs))
            if txt:
                out.append(txt)
    if getattr(shape, "has_table", False):
        tbl = shape.table
        for r in tbl.rows:
            cells = [_squeeze(c.text) for c in r.cells]
            if any(cells):
                out.append(" | ".join(c for c in cells if c))
    if getattr(shape, "shape_type", None) and str(shape.shape_type) == "GROUP":
        try:
            for sh in getattr(shape, "shapes", []):
                out.extend(_shape_text(sh))
        except Exception:
            pass
    try:
        S = SETTINGS.effective
        if bool(S().get("pptx_ocr_images", False)):
            is_pic = getattr(shape, "shape_type", None)
            if is_pic and "PICTURE" in str(is_pic):
                img = getattr(shape, "image", None)
                blob = getattr(img, "blob", None) if img is not None else None
                if blob and len(blob) >= int(S().get("ocr_min_image_bytes", 16384)):
                    log.info("[OCR] candidate image size: %d", len(blob))
                    try:
                        t = (ocr_image_bytes(blob) or "").strip()
                    except MissingTesseractError as e:
                        hint = _install_hint()
                        raise HTTPException(
                            status_code=424,
                            detail={
                                "code": "TESSERACT_MISSING",
                                "message": str(e),
                                "installUrl": hint["url"],
                                "note": hint["note"],
                            },
                        )
                    log.info("[OCR] result: %r", t[:200])
                    if t:
                        out.append(t)
    except HTTPException:
        # rethrow the 424 so FastAPI can serialize it
        raise
    except Exception as e:
        log.error("[OCR] error: %r", e)
    return out


def extract_pptx(data: bytes) -> tuple[str, str]:
    from pptx import Presentation

    S = SETTINGS.effective
    USE_MD = bool(S().get("pptx_use_markdown_headings", True))
    INCLUDE_NOTES = bool(S().get("pptx_include_notes", True))
    INCLUDE_TABLES = bool(S().get("pptx_include_tables", True))
    DROP_EMPTY = bool(S().get("pptx_drop_empty_lines", True))
    MAX_PARA = int(S().get("pptx_para_max_chars", 0))
    NUMBER_SLIDES = bool(S().get("pptx_number_slides", True))
    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        try:
            if getattr(slide, "shapes", None):
                for sh in slide.shapes:
                    if getattr(sh, "is_placeholder", False) and str(
                        getattr(sh, "placeholder_format", "").type
                    ).lower().endswith("title"):
                        title = _squeeze(getattr(sh, "text", "") or "")
                        break
        except Exception:
            pass
        head = f"Slide {i}" + (f": {title}" if title else "")
        if USE_MD:
            lines.append(("## " if NUMBER_SLIDES else "## ") + head)
        else:
            lines.append(head)
        body: list[str] = []
        for sh in getattr(slide, "shapes", []):
            if getattr(sh, "has_table", False) and (not INCLUDE_TABLES):
                continue
            body.extend(_shape_text(sh))
        for t in body:
            t = _clip(t, MAX_PARA)
            if t or not DROP_EMPTY:
                lines.append(t)
        if INCLUDE_NOTES:
            try:
                notes = slide.notes_slide
                if notes and getattr(notes, "notes_text_frame", None):
                    note_txt = _squeeze(notes.notes_text_frame.text)
                    if note_txt:
                        lines.append("")
                        lines.append("### Notes")
                        for ln in note_txt.splitlines():
                            ln = _squeeze(ln)
                            if ln or not DROP_EMPTY:
                                lines.append(_clip(ln, MAX_PARA))
            except Exception:
                pass
        lines.append("")
    text = "\n".join(line.rstrip() for line in lines if line is not None).strip()
    return (text + "\n" if text else "", "text/plain")


def extract_ppt(data: bytes) -> tuple[str, str]:
    from .doc_binary_ingest import _generic_ole_text

    S = SETTINGS.effective
    DROP_EMPTY = bool(S().get("ppt_drop_empty_lines", True))
    DEDUPE = bool(S().get("ppt_dedupe_lines", True))
    MAX_PARA = int(S().get("ppt_max_line_chars", 600))
    MIN_ALPHA = float(S().get("ppt_min_alpha_ratio", 0.4))
    MAX_PUNCT = float(S().get("ppt_max_punct_ratio", 0.5))
    TOKEN_MAX = int(S().get("ppt_token_max_chars", 40))
    raw = _generic_ole_text(data)
    if not raw:
        try:
            raw = data.decode("utf-8", errors="ignore")
        except Exception:
            raw = ""
    out: list[str] = []
    seen = set()
    for ln in raw.splitlines() if raw else []:
        s = _squeeze(ln)
        if not s and DROP_EMPTY:
            continue
        if MAX_PARA > 0 and len(s) > MAX_PARA:
            s = s[:MAX_PARA] + "…"
        if s:
            letters = sum(1 for c in s if c.isalpha())
            alen = max(1, len(s))
            if letters / alen < MIN_ALPHA:
                continue
            punct = sum(1 for c in s if not c.isalnum() and (not c.isspace()))
            if punct / alen > MAX_PUNCT:
                continue
            if " " not in s and len(s) <= TOKEN_MAX and re.fullmatch("[\\w.\\-]+", s):
                continue
        if DEDUPE:
            if s in seen:
                continue
            seen.add(s)
        if s or not DROP_EMPTY:
            out.append(s)
    text = "\n".join(out).strip()
    return (text + "\n" if text else "", "text/plain")
